import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


# Create PPT generator class
class PPTGenerator:
    # Initialize function, define PPT properties
    def __init__(self, title, slide_width=Inches(16), slide_height=Inches(9)):
        self.title = title
        self.prs = Presentation()
        self.slides = self.prs.slides
        # self.prs.slide_width = slide_width
        # self.prs.slide_height = slide_height

    def get_slide(self, index):
        return self.slides[index]

    def get_slide_layout(self, index):
        return self.prs.slide_layouts[index]

    def add_slide(self, layout):
        return self.slides.add_slide(layout)

    def delete_slide(self, index):
        self.slides._sldIdLst.remove(self.slides[index]._element)

    # Add title (and subtitle)
    def add_title(self, title_text, index_slide, subtitle=None):
        # title_slide_layout = self.prs.slide_layouts[0]
        # slide = self.prs.slides.add_slide(title_slide_layout)
        # title = slide.shapes.title
        # subtitle_box = slide.placeholders[1]
        #
        # title.text = title_text
        # subtitle.text = subtitle
        #
        slide = self.get_slide(index_slide)
        title_box = slide.shapes.title
        # title_box.margin_top = Inches(5)
        # title_box.margin_left = Inches(5)
        title_box.text = title_text
        subtitle_box = slide.placeholders[1]
        # title_box.margin_bottom = Inches()

        # title_box.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # title_box.word_wrap = False
        # title_box.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        if subtitle != None:
            subtitle_box.text = subtitle

    # Add Subtitle
    def add_subtitle(self, subtitle, index_slide, subtitle_left, subtitle_top):
        # slide = self.prs.slides[index_slide]
        # slide.shapes.title.text = subtitle

        slide = self.get_slide(index_slide)
        subtitle = slide.shapes.add_textbox(subtitle_left, subtitle_top, width=Inches(8), height=Inches(1))
        shape = slide.shapes
        title_shape = shape.title
        # title_shape.text_frame.paragraphs[0].font.size = Pt(10)
        title_shape.text_frame.paragraphs[0].font.size = Pt(15)
        # subtitle.text_frame.paragraphs[0].font.size = Pt(16)
        title_shape.text = subtitle

    # Add text
    def add_text(self, title, content, index_slide, location='left'):
        slide = self.prs.slides[index_slide]

        title_box = slide.shapes.title
        title_box.text = title
        # title_box.left = Inches(0.5)
        # title_box.top = Inches(0.3)
        # title_box.width = Inches(5)
        # title_box.height = Inches(1)
        if location == 'center':
            title_box.left = Inches(2.8)
            title_box.top = Inches(0.3)
            title_box.width = Inches(4.5)
            title_box.height = Inches(1)

        elif location == 'left':
            title_box.width = Inches(4.5)
            title_box.left = Inches(0.5)
            title_box.top = Inches(0.3)
            title_box.height = Inches(1)

        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.name = "Calibri"
        title_para.font.size = Pt(18)
        title_para.font.underline = True

        body_shape = slide.shapes.placeholders[1]
        if location == 'center':
            body_shape.left = Inches(1.5)
            body_shape.top = Inches(1.5)
            body_shape.width = Inches(8)
            body_shape.height = Inches(8)

        elif location == 'left':
            body_shape.width = Inches(5.0)
            body_shape.left = Inches(0.0)
            body_shape.top = Inches(1.5)
            body_shape.height = Inches(8)


        if type(content) == list:
            for i in range(len(content)):

                p = body_shape.text_frame.add_paragraph()
                p.text = content[i]

                text_place = body_shape.text_frame.paragraphs[i+1]
                text_place.font.name = "Calibri"
                text_place.font.size = Pt(14)
                # text_place.text = content[i]
                # text_place.width = Inches(5.5)
                # text_place.height = Inches(3)
                # text_place.margin_bottom = Inches(0.08)
                # text_place.margin_left = 0
                text_place.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                text_place.word_wrap = False
                text_place.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        else:
            text_place = body_shape.text_frame.paragraphs[0]
            text_place.font.name = "Calibri"
            text_place.font.size = Pt(14)
            text_place.text = content
            # text_place.width = Inches(5.5)
            # text_place.height = Inches(3)
            # text_place.margin_bottom = Inches(0.08)
            # text_place.margin_left = 0
            text_place.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            text_place.word_wrap = False
            text_place.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            text_place.alignment = PP_PARAGRAPH_ALIGNMENT.JUSTIFY


        # left = top = Inches(2)
        # width = height = Inches(6)
        # text_place = slide.shapes.add_textbox(left, top, width, height)
        # tf = text_place
        # tf.text = content
        # p = tf.add_paragraph()
        # p.font.size = Pt(14)
        # p.font.name = "Calibri"

        # text_place.font.name = "Calibri"
        # text_place.font.size = Pt(14)

    # Add image
    def add_image(self, title, image_file, index_slide):
        pic_left = Inches(5.25)
        pic_top = Inches(0.5)
        pic_width = Inches(4.5)
        slide = self.get_slide(index_slide)
        pic_height = self.prs.slide_height - Inches(1.0)
        img = slide.shapes.add_picture(image_file, pic_left, pic_top, pic_width, pic_height)

    # Set background

    # Save PPT
    def save_ppt(self):
        ppt_name = '{}.pptx'.format(self.title)
        self.prs.save(ppt_name)

